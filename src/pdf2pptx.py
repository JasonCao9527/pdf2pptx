# coding: UTF-8
# 本代码fork自https://github.com/phasedOut/pdf2pptx
# slide.shapes.add_picture的使用参考[python-pptx 0.6.21 documentation](https://python-pptx.readthedocs.io/en/latest/api/shapes.html)的指引
# 原文：
# The picture is positioned with its top-left corner at (top, left). If width and height are both , the native size of the image is used.(换行)
# If only one of width or height is used, the unspecified dimension is calculated to preserve the aspect ratio of the image.
# 删改说明：
# * 删除了使得生成的pptx文件中插入生成时间的代码，让功能更加纯粹
# * 调整了原来的slide.shapes.add_picture插入的位置，现在不会「跑偏」

#python 报错 AttributeError: module 'collections' has no attribute 'Container'的解决方法如下：
#这个错误是由于Python 3.10中的collections模块中已经删除了MutableMapping类，因此您会看到AttributeError，
#其中写着“module ‘collections’ has no attribute ‘MutableMapping’”。
#collections模块提供了各种容器数据类型，但是在Python 3.10及更高版本中，MutableMapping类已被删除。
#解决方法是在导入pptx之前导入collections和collections.abc两个模块
import collections 
import collections.abc 

from pptx import Presentation
from pdf2image import convert_from_path
import os
from PIL import Image

for filename in os.listdir('source_files/'):
    if os.path.splitext(filename)[1] == '.pdf':
        print("Creating %s" % filename)
        #prs = Presentation()
        # 下面代码是为了解决打包成exe之后pptx库的报错
        # 声明一个pptx对象时加入如下参数
        #不加下面的参数打包成exe之后就会报pptx.exc.PackageNotFoundError
        #错误内容：pptx.exc.PackageNotFoundError: Package not found at 'C:\Users\29571\AppData\Local\Temp\_MEI427362\pptx\templates\default.pptx'
        #解决方法来自：https://blog.csdn.net/weixin_54693379/article/details/128072858
        #n.b.
        #同时，需要将安装的pptx包中的default.pptx复制到exe文件同目录，default.pptx在pptx / template文件夹下
        #同理，docx库报这个错误也是这种解决办法
        #本程序中还需要把三个文件夹复制到exe的同目录下，才可以正常工作
        #
        prs = Presentation(pptx=os.path.join(os.getcwd(), 'default.pptx'))

        pages = convert_from_path('source_files/' + filename, 700)
        # 这行代码的作用是将名为 filename 的文件从路径 source_files/ 中读取，然后将其转换为一个包含多个 JPEG 图像的列表，
        # 每个图像都是从 PDF 文件的一个页面中提取的。这里的 500 是指每英寸像素数（PPI），它控制了输出图像的分辨率。
        for index, page in enumerate(pages):
            #Save as 'jpg' in jpgs dir
            jpg_file = "jpgs/%s-(%d).jpg" % (filename,index)
            page.save(jpg_file, 'JPEG')

            #Get width/height of image
            image = Image.open(jpg_file)
            height = image.height
            width = image.width
            #Rotate 270 degrees if horizontal
            if height > width:
                adjusted = image.rotate(270, expand=True)
                adjusted.save(jpg_file)

            #Setup slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            left = top = 0
            # slide.shapes.add_picture(jpg_file, left-0.1*prs.slide_width,top,height = prs.slide_height)
            slide.shapes.add_picture(jpg_file, left,top,height = prs.slide_height)
        prs.save('result/%s.pptx' % os.path.splitext(filename)[0])

    else:
        print("Skipping %s because it\'s not a pdf" % filename)

print("Saved to result directory")
os.system("pause")
